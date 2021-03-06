from pyate.term_extraction_pipeline import TermExtractionPipeline
from nltk.tokenize import word_tokenize
from stop_words import get_stop_words
from email.parser import BytesParser
from openpyxl import load_workbook
from odf.opendocument import load
from odf import text, teletype
from pptx import Presentation
from bs4 import BeautifulSoup
from docx2txt import docx2txt
from pprint import pprint
from email import policy
import spacy_udpipe
import extract_msg
import textract
import PyPDF2
import glob
import re
import os

PATH = "test.html"
PATH_MODELS = "models/czech-pdt-ud-2.5-191206.udpipe"


class Extract:
    def __init__(self, path: str):
        self.UNIT_TO_MULTIPLIER = {

            '.csv': self.get_text_with_txt,
            '.doc': self.get_text_with_docx,
            '.docx': self.get_text_with_docx,
            '.eml': self.get_text_with_eml,
            '.html': self.get_text_with_html,
            '.msg': self.get_text_with_msg,
            '.pdf': self.get_text_with_pdf,
            '.pptx': self.get_text_with_pptx,
            '.txt': self.get_text_with_txt,
            '.xlsx': self.get_text_with_xml,
            '.xls': self.get_text_with_xml,
            '.idml': None,
            '.mht': None,
            '.odt': self.get_text_with_odt,
            '.odp': None,
            '.ods': None,
            '.xml': None,
            '.xps': None,
            '.zip': None
        }
        self.path = path

    def extract(self):
        return self.switch_case()

    def get_format(self):
        filename, file_extension = os.path.splitext(self.path)
        return file_extension

    def switch_case(self) -> str:
        try:

            command = self.UNIT_TO_MULTIPLIER[self.get_format()]
            return command() if command is not None else "This file format is not supported"

        except Exception as e:
            return e

    def get_text_with_eml(self) -> str:

        file_list = glob.glob('*.eml')  # returns list of files
        with open(file_list[2], 'rb') as fp:  # select a specific email file from the list
            msg = BytesParser(policy=policy.default).parse(fp)
        return msg.get_body(preferencelist=('plain')).get_content()

    def get_text_with_xml(self) -> str:
        workbook = load_workbook(self.path)
        first_sheet = workbook.get_sheet_names()[0]
        worksheet = workbook.get_sheet_by_name(first_sheet)
        text = ""
        for row in worksheet.iter_rows():
            text += str(row)

        # check out the last row
        for cell in row:
            text += str(cell)
        return text

    def get_text_with_docx(self) -> str:
        return docx2txt.process(self.path)

    def get_text_with_pptx(self):
        result = ""
        for eachfile in glob.glob(self.path):
            prs = Presentation(eachfile)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += shape.text
        return result

    def get_text_with_txt(self) -> str:
        return textract.process(self.path)

    def get_text_with_html(self) -> str:
        with open(self.path) as fp:
            soup = BeautifulSoup(fp, features="html.parser")

        # kill all script and style elements
        for script in soup(["script", "style"]):
            script.extract()  # rip it out

        # get text
        text = soup.get_text()

        # break into lines and remove leading and trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return text

    def get_text_with_msg(self) -> str:
        return extract_msg.openMsg(self.path)

    def get_text_with_odt(self) -> str:
        textdoc = load(self.path)
        allparas = textdoc.getElementsByType(text.P)
        return teletype.extractText(allparas)

    def get_text_with_pdf(self) -> str:
        # creating a pdf file object
        pdfFileObj = open('text/test.pdf', 'rb')

        # creating a pdf reader object
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        text = ""
        for i in range(int(pdfReader.numPages)):
            # creating a page object
            pageObj = pdfReader.getPage(i)

            # extracting text from page
            text += pageObj.extractText()

        # closing the pdf file object
        pdfFileObj.close()
        return text


def remove_all_except_letter_dot(text: str) -> str:
    return re.sub(r"[^AaÁáBbCcČčDdĎďEeÉéĔĕFfGgHhСhchIiÍíJjKkLlMmNnŇňOoÓóPpQqRrŘřSsŠšTtŤťUuÚúVvWwXxYyÝýZzŽž.]", " ",
                  text)


def remove_stop_words_cs(text: str) -> str:
    stop_words = set(get_stop_words('czech'))
    word_tokens = word_tokenize(text)

    filtered_sentence = []

    for w in word_tokens:
        if w not in stop_words and len(w) > 3:
            filtered_sentence.append(w)

    return " ".join(filtered_sentence)


ex = Extract(PATH)
# pprint(ex.extract())
print("=" * 100)
text = remove_all_except_letter_dot(ex.extract())
text = remove_stop_words_cs(text=text)

nlp = spacy_udpipe.load_from_path(lang="cs",
                                  path=PATH_MODELS,
                                  meta={"description": "Custom 'hr' model"})

nlp.add_pipe(TermExtractionPipeline())

doc = nlp(text)
pprint(list(dict(doc._.combo_basic.sort_values(ascending=False)).keys()))
